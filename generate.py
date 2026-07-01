"""
Generate a styled PowerPoint presentation from the LASU Certificate Generator HTML content
This script creates a fully formatted PPTX with all slides, content, and image placeholders
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Color scheme matching the HTML theme
class Colors:
    BG = RGBColor(26, 24, 23)           # #1a1817
    TEXT = RGBColor(250, 249, 245)      # #faf9f5
    ACCENT = RGBColor(217, 119, 87)     # #d97757
    SUBTEXT = RGBColor(176, 174, 165)   # #b0aea5
    MUTED = RGBColor(106, 106, 101)     # #6a6a65
    CARD_BG = RGBColor(42, 37, 34)      # #2a2522
    BORDER = RGBColor(58, 53, 50)       # #3a3532
    BLUE = RGBColor(106, 155, 204)      # #6a9bcc
    GREEN = RGBColor(120, 140, 93)      # #788c5d

def set_slide_background(slide, color=Colors.BG):
    """Set slide background color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, text, left, top, width, height, 
                 font_size=18, color=Colors.TEXT, bold=False, 
                 italic=False, align=PP_ALIGN.LEFT, font_name='Arial'):
    """Helper to add formatted text box"""
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = text
    text_frame.paragraphs[0].font.size = Pt(font_size)
    text_frame.paragraphs[0].font.color.rgb = color
    text_frame.paragraphs[0].font.bold = bold
    text_frame.paragraphs[0].font.italic = italic
    text_frame.paragraphs[0].font.name = font_name
    text_frame.paragraphs[0].alignment = align
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    return text_box

def add_rounded_rect(slide, left, top, width, height, 
                     fill_color=Colors.CARD_BG, 
                     border_color=Colors.BORDER,
                     border_width=Pt(1)):
    """Add a rounded rectangle shape"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = border_width
    return shape

def add_tag(slide, text, left, top, color=Colors.ACCENT):
    """Add a styled tag/badge"""
    tag_width = Inches(len(text) * 0.12 + 0.4)
    tag_height = Inches(0.35)
    
    # Tag background (transparent with border)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, tag_width, tag_height
    )
    shape.fill.background()
    shape.line.color.rgb = color
    shape.line.width = Pt(1)
    
    # Tag text
    text_box = slide.shapes.add_textbox(
        left + Inches(0.1), top + Inches(0.02),
        tag_width - Inches(0.2), tag_height - Inches(0.04)
    )
    text_frame = text_box.text_frame
    text_frame.text = text
    text_frame.paragraphs[0].font.size = Pt(10)
    text_frame.paragraphs[0].font.color.rgb = color
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def create_presentation():
    """Create the full presentation"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ==================== SLIDE 1: TITLE ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    
    # Main title
    add_text_box(slide, "LASU Certificate", 
                 Inches(0.8), Inches(1.5), Inches(11), Inches(1.2),
                 font_size=48, color=Colors.ACCENT, bold=True)
    add_text_box(slide, "Generator", 
                 Inches(0.8), Inches(2.5), Inches(11), Inches(1.0),
                 font_size=48, color=Colors.ACCENT, bold=True)
    
    # Subtitle
    add_text_box(slide, "OOP Design Patterns • C++17 • Winsock HTTP Server • Web UI",
                 Inches(0.8), Inches(3.8), Inches(11), Inches(0.6),
                 font_size=20, color=Colors.SUBTEXT)
    
    # Course info
    add_text_box(slide, "CSC 226 — Object-Oriented Programming II — Lagos State University — July 2026",
                 Inches(0.8), Inches(4.5), Inches(11), Inches(0.5),
                 font_size=14, color=Colors.TEXT, bold=True)
    
    # Tags
    tags = ["Strategy", "Template Method", "Factory", "Inheritance", 
            "Polymorphism", "RAII", "std::map", "Multithreading"]
    tag_colors = [Colors.ACCENT, Colors.BLUE, Colors.GREEN, Colors.ACCENT,
                  Colors.BLUE, Colors.GREEN, Colors.ACCENT, Colors.BLUE]
    x_pos = Inches(0.8)
    for tag, color in zip(tags, tag_colors):
        add_tag(slide, tag, x_pos, Inches(5.1), color)
        x_pos += Inches(len(tag) * 0.12 + 0.5)
    
    # Team
    team_text = ("Group 6: Idowu Fawaz, Idowu Matthew, Ikechukwu Queen, Inedu Esther, Iniotuh Greatest, "
                 "Iwakun Oluwasegun, Iwuno Vincent, Iyaniwura George, Jawando Fuad, Jeremiah David, "
                 "Joseph Elizabeth, Kalejaiye Halimah, Kasali Damilola, Kassim Oluwatimilehin, "
                 "Kehinde Oyindamola, Kolawole Abubakar, Kwegan Sean, Lamina Rihanat, Lawal Sahal, Lawal Ibrahim")
    add_text_box(slide, team_text,
                 Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.8),
                 font_size=9, color=Colors.MUTED, align=PP_ALIGN.CENTER)
    
    # ==================== SLIDE 2: OVERVIEW & ARCHITECTURE ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    
    add_text_box(slide, "Overview & Architecture",
                 Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                 font_size=36, color=Colors.ACCENT, bold=True)
    
    add_text_box(slide, "Desktop web app: generate, manage, and print HTML certificates",
                 Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
                 font_size=18, color=Colors.SUBTEXT, italic=True)
    
    # Feature cards
    features = [
        ("3 Certificate Types", "Excellence, Completion, Participation"),
        ("3 Template Styles", "Formal (serif), Modern (gradient), Minimalist (mono)"),
        ("CSV Batch Import", "File path or paste • 10 sample students"),
        ("Search, Sort & Stats", "Name, course, GPA • Per-course & grade distro")
    ]
    
    for i, (title, desc) in enumerate(features):
        row = i // 2
        col = i % 2
        x = Inches(0.8 + col * 5.8)
        y = Inches(1.8 + row * 1.2)
        
        add_rounded_rect(slide, x, y, Inches(5.3), Inches(0.9))
        add_text_box(slide, title,
                     x + Inches(0.15), y + Inches(0.05),
                     Inches(5), Inches(0.4),
                     font_size=14, color=Colors.TEXT, bold=True)
        add_text_box(slide, desc,
                     x + Inches(0.15), y + Inches(0.4),
                     Inches(5), Inches(0.4),
                     font_size=11, color=Colors.SUBTEXT)
    
    # Architecture diagram
    arch_y = Inches(4.2)
    arch_items = ["Browser", "HttpServer\nWinsock TCP", "Certificate\nGenerator\nstd::map CRUD", "Certificate\nTemplate Method"]
    arch_colors = [Colors.BLUE, Colors.BLUE, Colors.ACCENT, Colors.GREEN]
    item_width = Inches(2.4)
    spacing = Inches(0.3)
    total_width = len(arch_items) * (item_width + spacing) - spacing
    start_x = (prs.slide_width - total_width) / 2
    
    for i, (text, color) in enumerate(zip(arch_items, arch_colors)):
        x = start_x + i * (item_width + spacing)
        
        # Box
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, arch_y, item_width, Inches(0.9)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = Colors.CARD_BG
        shape.line.color.rgb = color
        shape.line.width = Pt(2)
        
        # Text
        add_text_box(slide, text,
                     x + Inches(0.1), arch_y + Inches(0.05),
                     item_width - Inches(0.2), Inches(0.8),
                     font_size=12, color=Colors.TEXT, align=PP_ALIGN.CENTER)
        
        # Arrow
        if i < len(arch_items) - 1:
            add_text_box(slide, "→",
                         x + item_width, arch_y + Inches(0.2),
                         Inches(0.3), Inches(0.5),
                         font_size=24, color=Colors.ACCENT, align=PP_ALIGN.CENTER)
    
    # Data flow
    add_text_box(slide, "Data flow: Browser → HTTP (JSON) → HttpServer → CertificateGenerator → Certificate + Template → HTML file",
                 Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.4),
                 font_size=11, color=Colors.MUTED, align=PP_ALIGN.CENTER)
    
    # ==================== SLIDE 3: INHERITANCE & ENCAPSULATION ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    
    add_text_box(slide, "Inheritance & Encapsulation",
                 Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                 font_size=36, color=Colors.ACCENT, bold=True)
    
    add_text_box(slide, "Person (abstract) → Student (concrete) • All data private",
                 Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
                 font_size=18, color=Colors.SUBTEXT, italic=True)
    
    # Code blocks
    code1 = '''class Person {                   // Abstract base
protected:
    std::string name_;
public:
    virtual ~Person() = default;
    virtual std::string
        getDescription() const = 0; // pure virtual
};

class Student : public Person {  // is-a Person
    std::string course_, grade_;
public:
    double gradeToPoints() const;
    // 'A'->4.0  'B+'->3.3  '93'->3.72
    static bool isValidGrade(const&);
    std::string getDescription()
        const override;
};'''
    
    code2 = '''// Encapsulation in Generator
class CertificateGenerator {
    std::map<std::string, Student>
        students_;              // private
    mutable std::mutex mutex_;  // private
public:
    void addStudent(const&);    // validated
    bool removeStudent(const&);
    Student* findStudent(const&);
    void loadFromCsv(const&);
};'''
    
    # Code block 1
    code_bg1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.9),
        Inches(5.5), Inches(3.2)
    )
    code_bg1.fill.solid()
    code_bg1.fill.fore_color.rgb = Colors.CARD_BG
    code_bg1.line.color.rgb = Colors.BORDER
    
    add_text_box(slide, code1,
                 Inches(1.0), Inches(2.0),
                 Inches(5.1), Inches(3.0),
                 font_size=9, color=Colors.TEXT, font_name='Courier New')
    
    # Code block 2
    code_bg2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.9),
        Inches(5.7), Inches(3.2)
    )
    code_bg2.fill.solid()
    code_bg2.fill.fore_color.rgb = Colors.CARD_BG
    code_bg2.line.color.rgb = Colors.BORDER
    
    add_text_box(slide, code2,
                 Inches(7.0), Inches(2.0),
                 Inches(5.3), Inches(3.0),
                 font_size=9, color=Colors.TEXT, font_name='Courier New')
    
    # Bullets
    add_text_box(slide, "• Inheritance: Student is-a Person — inherits name_, adds course/grade",
                 Inches(0.8), Inches(5.3), Inches(11.7), Inches(0.4),
                 font_size=14, color=Colors.SUBTEXT)
    add_text_box(slide, "• Encapsulation: All data private — addStudent() validates before storing, throws on bad data",
                 Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.4),
                 font_size=14, color=Colors.SUBTEXT)
    
    # ==================== SLIDE 4: POLYMORPHISM & ABCs ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    
    add_text_box(slide, "Polymorphism & Abstract Base Classes",
                 Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                 font_size=36, color=Colors.ACCENT, bold=True)
    
    add_text_box(slide, "3 ABCs, 7 pure virtuals, runtime vtable dispatch",
                 Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
                 font_size=18, color=Colors.SUBTEXT, italic=True)
    
    # Code block
    code3 = '''class Certificate {              // Abstract base
public:
    virtual std::string
        generateHTML() const;    // Template Method
    virtual std::string
        getTypeName() const = 0; // pure virtual
};

class CertOfExcellence : public Certificate {
    std::string getTypeName() const override
        { return "Excellence"; }
};

class CertOfCompletion : public Certificate {
    std::string getTypeName() const override
        { return "Completion"; }
};'''
    
    code_bg3 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8),
        Inches(5.5), Inches(3.5)
    )
    code_bg3.fill.solid()
    code_bg3.fill.fore_color.rgb = Colors.CARD_BG
    code_bg3.line.color.rgb = Colors.BORDER
    
    add_text_box(slide, code3,
                 Inches(1.0), Inches(1.9),
                 Inches(5.1), Inches(3.3),
                 font_size=9, color=Colors.TEXT, font_name='Courier New')
    
    # Key points
    points = [
        "3 ABCs: Person, CertificateTemplate, Certificate",
        "7 pure virtual methods total",
        "Runtime dispatch via vtable",
        "Client code works with base pointers"
    ]
    
    y_pos = Inches(1.9)
    for point in points:
        add_text_box(slide, f"• {point}",
                     Inches(6.8), y_pos,
                     Inches(5.7), Inches(0.4),
                     font_size=14, color=Colors.SUBTEXT)
        y_pos += Inches(0.5)
    
    # Code snippet
    code4 = '''// Factory returns subclass at runtime
auto tmpl = makeTemplate(style);
auto cert = makeCertificate(
    type, s, std::move(tmpl));
file << cert->generateHTML();   // polymorphic!'''
    
    code_bg4 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(4.0),
        Inches(5.7), Inches(1.3)
    )
    code_bg4.fill.solid()
    code_bg4.fill.fore_color.rgb = Colors.CARD_BG
    code_bg4.line.color.rgb = Colors.BORDER
    
    add_text_box(slide, code4,
                 Inches(7.0), Inches(4.1),
                 Inches(5.3), Inches(1.1),
                 font_size=9, color=Colors.TEXT, font_name='Courier New')
    
    # ==================== SLIDE 5: STRATEGY + TEMPLATE METHOD ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    
    add_text_box(slide, "Strategy + Template Method",
                 Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                 font_size=36, color=Colors.ACCENT, bold=True)
    
    add_text_box(slide, "Strategy: swappable visual styles • Template Method: fixed skeleton",
                 Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
                 font_size=18, color=Colors.SUBTEXT, italic=True)
    
    # Code block 1
    code5 = '''// STRATEGY: CertificateTemplate
class CertificateTemplate {
public:
    virtual string header(...) const = 0;
    virtual string styles() const = 0;
    virtual string body(...) const = 0;
};
class FormalTemplate : public CertTempl {
    // Double border, serif (Georgia)
};
class ModernTemplate : public CertTempl {
    // Gradient bg, rounded corners
};
class MinimalistTemplate : public CertTempl {
    // Monospace font, clean lines
};'''
    
    code_bg5 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8),
        Inches(5.5), Inches(3.2)
    )
    code_bg5.fill.solid()
    code_bg5.fill.fore_color.rgb = Colors.CARD_BG
    code_bg5.line.color.rgb = Colors.BORDER
    
    add_text_box(slide, code5,
                 Inches(1.0), Inches(1.9),
                 Inches(5.1), Inches(3.0),
                 font_size=8, color=Colors.TEXT, font_name='Courier New')
    
    # Code block 2
    code6 = '''// TEMPLATE METHOD: generateHTML()
string Certificate::generateHTML() const {
    string typeName = getTypeName(); // virtual
    string html;
    html += template_.header(typeName);
    html += template_.styles();
    html += template_.body(student_,
                           typeName, date);
    html += "</body></html>";
    return html;
}
// Subclass provides type name/desc
// Template provides all HTML/CSS'''
    
    code_bg6 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8),
        Inches(5.7), Inches(3.2)
    )
    code_bg6.fill.solid()
    code_bg6.fill.fore_color.rgb = Colors.CARD_BG
    code_bg6.line.color.rgb = Colors.BORDER
    
    add_text_box(slide, code6,
                 Inches(7.0), Inches(1.9),
                 Inches(5.3), Inches(3.0),
                 font_size=8, color=Colors.TEXT, font_name='Courier New')
    
    # Bullets
    add_text_box(slide, "• Strategy: Open/Closed — new templates don't change Certificate | Style chosen at runtime",
                 Inches(0.8), Inches(5.3), Inches(11.7), Inches(0.4),
                 font_size=13, color=Colors.SUBTEXT)
    add_text_box(slide, "• Template Method: getTypeName() dispatches polymorphically | Zero duplication across types",
                 Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.4),
                 font_size=13, color=Colors.SUBTEXT)
    
    # ==================== SLIDE 6: FACTORY + COMPOSITION ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    
    add_text_box(slide, "Factory + Composition",
                 Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                 font_size=36, color=Colors.ACCENT, bold=True)
    
    add_text_box(slide, "makeTemplate/makeCertificate • has-a relationships throughout",
                 Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
                 font_size=18, color=Colors.SUBTEXT, italic=True)
    
    # Code block 1
    code7 = '''// Factory: returns correct subclass
unique_ptr<CertificateTemplate>
makeTemplate(TemplateStyle style) {
    switch (style) {
    case Formal:
        return mk<FormalTemplate>();
    case Modern:
        return mk<ModernTemplate>();
    case Minimalist:
        return mk<MinimalistTemplate>();
    }
}'''
    
    code_bg7 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8),
        Inches(5.5), Inches(2.2)
    )
    code_bg7.fill.solid()
    code_bg7.fill.fore_color.rgb = Colors.CARD_BG
    code_bg7.line.color.rgb = Colors.BORDER
    
    add_text_box(slide, code7,
                 Inches(1.0), Inches(1.9),
                 Inches(5.1), Inches(2.0),
                 font_size=8, color=Colors.TEXT, font_name='Courier New')
    
    # Code block 2
    code8 = '''// Composition: has-a
class Certificate {
    const Student& student_;
    unique_ptr<CertificateTemplate>
        template_; // has-a Template
};'''
    
    code_bg8 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.2),
        Inches(5.5), Inches(1.3)
    )
    code_bg8.fill.solid()
    code_bg8.fill.fore_color.rgb = Colors.CARD_BG
    code_bg8.line.color.rgb = Colors.BORDER
    
    add_text_box(slide, code8,
                 Inches(1.0), Inches(4.3),
                 Inches(5.1), Inches(1.1),
                 font_size=8, color=Colors.TEXT, font_name='Courier New')
    
    # Key points
    points = [
        "Factory: Specify an enum → get right polymorphic type",
        "No conditional logic at call site",
        "Add type: enum value + case + class",
        "Composition: has-a beats is-a",
        "Runtime flexibility (swap template)",
        "Looser coupling, easier testing",
        "HttpServer has-a Generator",
        "Generator has-a map of Students"
    ]
    
    y_pos = Inches(1.9)
    for point in points:
        add_text_box(slide, f"• {point}",
                     Inches(6.8), y_pos,
                     Inches(5.7), Inches(0.35),
                     font_size=11, color=Colors.SUBTEXT)
        y_pos += Inches(0.4)
    
    # ==================== SLIDE 7: RAII + STL ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    
    add_text_box(slide, "RAII & STL Features",
                 Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                 font_size=36, color=Colors.ACCENT, bold=True)
    
    add_text_box(slide, "Automatic resource management • std::map, sort, find_if, lambdas, thread",
                 Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
                 font_size=18, color=Colors.SUBTEXT, italic=True)
    
    # Code block 1
    code9 = '''// RAII: no manual cleanup
void addStudent(const Student& s) {
    std::lock_guard<mutex> lock(m_);
    students_[toLower(s.getName())] = s;
    // auto-unlocked at scope exit
}

// unique_ptr: auto-deleted polymorphic
auto cert = makeCertificate(type, s, t);
// cert + tmpl auto-deleted when done'''
    
    code_bg9 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8),
        Inches(5.5), Inches(3.0)
    )
    code_bg9.fill.solid()
    code_bg9.fill.fore_color.rgb = Colors.CARD_BG
    code_bg9.line.color.rgb = Colors.BORDER
    
    add_text_box(slide, code9,
                 Inches(1.0), Inches(1.9),
                 Inches(5.1), Inches(2.8),
                 font_size=8, color=Colors.TEXT, font_name='Courier New')
    
    # Code block 2
    code10 = '''// STL: case-insensitive search
Student* findStudent(const string& n) {
    auto it = find_if(begin(students_),
        end(students_),
        [&](const auto& p) {
            return toLower(p.first)
                   == toLower(n);
        });
    return it != end(students_)
           ? &it->second : nullptr;
}'''
    
    code_bg10 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8),
        Inches(5.7), Inches(2.5)
    )
    code_bg10.fill.solid()
    code_bg10.fill.fore_color.rgb = Colors.CARD_BG
    code_bg10.line.color.rgb = Colors.BORDER
    
    add_text_box(slide, code10,
                 Inches(7.0), Inches(1.9),
                 Inches(5.3), Inches(2.3),
                 font_size=8, color=Colors.TEXT, font_name='Courier New')
    
    # Bullets
    add_text_box(slide, "• RAII: unique_ptr, lock_guard, ofstream — zero manual delete/unlock/close",
                 Inches(0.8), Inches(5.3), Inches(11.7), Inches(0.4),
                 font_size=13, color=Colors.SUBTEXT)
    add_text_box(slide, "• STL: std::map (O(log n)), find_if+lambda, sort+lambda, thread, mutex, stringstream",
                 Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.4),
                 font_size=13, color=Colors.SUBTEXT)
    
    # ==================== SLIDE 8: HTTP SERVER ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    
    add_text_box(slide, "HTTP Server & REST API",
                 Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                 font_size=36, color=Colors.ACCENT, bold=True)
    
    add_text_box(slide, "Winsock TCP • 12 endpoints • Embedded SPA • Thread safety",
                 Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
                 font_size=18, color=Colors.SUBTEXT, italic=True)
    
    # Code block
    code11 = '''// Accept loop: each conn on its own thread
while (true) {
    SOCKET client = accept(server,
        nullptr, nullptr);
    if (client == INVALID_SOCKET) break;
    std::thread(handleClient, client,
        std::ref(gen)).detach();
}'''
    
    code_bg11 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8),
        Inches(5.5), Inches(2.0)
    )
    code_bg11.fill.solid()
    code_bg11.fill.fore_color.rgb = Colors.CARD_BG
    code_bg11.line.color.rgb = Colors.BORDER
    
    add_text_box(slide, code11,
                 Inches(1.0), Inches(1.9),
                 Inches(5.1), Inches(1.8),
                 font_size=8, color=Colors.TEXT, font_name='Courier New')
    
    # Features
    features = [
        "All HTML/CSS/JS in C++ raw string literal",
        "Frontend uses fetch() — no reloads",
        "Theme toggle: Light/Dark/System"
    ]
    y_pos = Inches(4.0)
    for feature in features:
        add_text_box(slide, f"• {feature}",
                     Inches(0.8), y_pos,
                     Inches(5.5), Inches(0.35),
                     font_size=12, color=Colors.SUBTEXT)
        y_pos += Inches(0.4)
    
    # API Table
    endpoints = [
        ("GET", "/", "Serve Web UI"),
        ("GET", "/list", "All students JSON"),
        ("GET", "/stats", "Statistics JSON"),
        ("POST", "/add", "Insert student"),
        ("POST", "/remove", "Delete student"),
        ("POST", "/load", "CSV import"),
        ("POST", "/generate", "Create certificates"),
        ("POST", "/sort", "Sort students")
    ]
    
    table_x = Inches(6.8)
    table_y = Inches(1.8)
    col_widths = [Inches(0.8), Inches(1.0), Inches(3.5)]
    
    # Table headers
    headers = ["Method", "Path", "Action"]
    x_pos = table_x
    for i, header in enumerate(headers):
        header_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x_pos, table_y, col_widths[i], Inches(0.4)
        )
        header_bg.fill.solid()
        header_bg.fill.fore_color.rgb = Colors.CARD_BG
        header_bg.line.color.rgb = Colors.BORDER
        add_text_box(slide, header,
                     x_pos + Inches(0.05), table_y + Inches(0.02),
                     col_widths[i] - Inches(0.1), Inches(0.35),
                     font_size=11, color=Colors.TEXT, bold=True, align=PP_ALIGN.CENTER)
        x_pos += col_widths[i]
    
    # Table rows
    for i, (method, path, action) in enumerate(endpoints):
        y_pos = table_y + Inches(0.4 + i * 0.35)
        x_pos = table_x
        row_data = [method, path, action]
        for j, data in enumerate(row_data):
            cell_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x_pos, y_pos, col_widths[j], Inches(0.35)
            )
            cell_bg.fill.solid()
            cell_bg.fill.fore_color.rgb = Colors.BG
            cell_bg.line.color.rgb = Colors.BORDER
            add_text_box(slide, data,
                         x_pos + Inches(0.05), y_pos + Inches(0.02),
                         col_widths[j] - Inches(0.1), Inches(0.3),
                         font_size=9, color=Colors.SUBTEXT, align=PP_ALIGN.CENTER)
            x_pos += col_widths[j]
    
    # Additional info
    add_text_box(slide, "• std::mutex guards map (thread-safe) • ValidationException → JSON error • CSV: invalid lines skipped",
                 Inches(6.8), Inches(5.8), Inches(5.7), Inches(0.4),
                 font_size=11, color=Colors.MUTED)
    
    # ==================== SLIDE 9: WEB UI PREVIEW ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    
    add_text_box(slide, "Web UI Preview",
                 Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                 font_size=36, color=Colors.ACCENT, bold=True)
    
    add_text_box(slide, "Browser-based interface — Dashboard, Students, Generate, Statistics tabs",
                 Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
                 font_size=18, color=Colors.SUBTEXT, italic=True)
    
    # Image placeholder
    img_placeholder = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8),
        Inches(7.5), Inches(4.5)
    )
    img_placeholder.fill.solid()
    img_placeholder.fill.fore_color.rgb = Colors.CARD_BG
    img_placeholder.line.color.rgb = Colors.BORDER
    img_placeholder.line.width = Pt(2)
    
    # Add placeholder text
    add_text_box(slide, "📸 [SCREENSHOT PLACEHOLDER]",
                 Inches(0.8), Inches(3.5),
                 Inches(7.5), Inches(0.6),
                 font_size=24, color=Colors.MUTED, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Web UI Screenshot",
                 Inches(0.8), Inches(4.2),
                 Inches(7.5), Inches(0.4),
                 font_size=14, color=Colors.SUBTEXT, align=PP_ALIGN.CENTER)
    
    # Features list
    features = [
        "Dashboard: Total students, avg GPA, cert count, active courses",
        "Students: Search, sort, add, remove — grade badges with colors",
        "Generate: Pick type + template style — one student or all",
        "Statistics: Course distribution, grade distribution",
        "Theme toggle: Light / Dark / Match System"
    ]
    
    y_pos = Inches(1.9)
    for feature in features:
        add_text_box(slide, f"• {feature}",
                     Inches(8.8), y_pos,
                     Inches(3.9), Inches(0.4),
                     font_size=11, color=Colors.SUBTEXT)
        y_pos += Inches(0.5)
    
    # ==================== SLIDE 10: ALL CONCEPTS IN ACTION ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    
    add_text_box(slide, "All Concepts in Action",
                 Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
                 font_size=32, color=Colors.ACCENT, bold=True)
    
    add_text_box(slide, "End-to-end: add student → generate certificate → HTML output",
                 Inches(0.8), Inches(0.9), Inches(11), Inches(0.4),
                 font_size=16, color=Colors.SUBTEXT, italic=True)
    
    # Code block 1
    code12 = '''// 1. Add (encapsulation + exception)
try {
    gen.addStudent(
        Student("Alice","CSC 226","A"));
} catch (const ValidationE& e) {
    cerr << e.what();
}

// 2. Create Template (factory + strategy)
auto tmpl = makeTemplate(Formal);

// 3. Create Certificate (factory + poly.)
auto cert = makeCertificate(
    Excellence, student, move(tmpl));'''
    
    code_bg12 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4),
        Inches(5.5), Inches(3.0)
    )
    code_bg12.fill.solid()
    code_bg12.fill.fore_color.rgb = Colors.CARD_BG
    code_bg12.line.color.rgb = Colors.BORDER
    
    add_text_box(slide, code12,
                 Inches(1.0), Inches(1.5),
                 Inches(5.1), Inches(2.8),
                 font_size=8, color=Colors.TEXT, font_name='Courier New')
    
    # Code block 2
    code13 = '''// 4. Generate HTML (template method)
string html = cert->generateHTML();
//   getTypeName() → "Excellence"
//   FormalTemplate::header() + styles()
//   FormalTemplate::body()

// 5. Write file (RAII: ofstream)
ofstream file("output/Alice.html");
file << html;  // auto-closed

// 6. Serve via HTTP (thread + mutex)
// Browser GET /list returns JSON'''
    
    code_bg13 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.4),
        Inches(5.7), Inches(3.0)
    )
    code_bg13.fill.solid()
    code_bg13.fill.fore_color.rgb = Colors.CARD_BG
    code_bg13.line.color.rgb = Colors.BORDER
    
    add_text_box(slide, code13,
                 Inches(7.0), Inches(1.5),
                 Inches(5.3), Inches(2.8),
                 font_size=8, color=Colors.TEXT, font_name='Courier New')
    
    # Concept cards
    concepts = [
        ("Inheritance", "Student : Person"),
        ("Polymorphism", "cert->generateHTML()"),
        ("Encapsulation", "private + getters"),
        ("Abstract Base", "3 ABCs, 7 virtuals"),
        ("Strategy", "FormalTemplate CSS"),
        ("Template Method", "generateHTML()"),
        ("Factory", "makeTemplate/Cert"),
        ("Composition", "has-a Template"),
        ("RAII", "unique_ptr/lock"),
        ("STL", "map/sort/thread"),
        ("Exception", "try/catch validate"),
        ("Thread Safety", "mutex guard")
    ]
    
    card_width = Inches(2.2)
    card_height = Inches(0.65)
    spacing = Inches(0.15)
    start_x = (prs.slide_width - (4 * (card_width + spacing) - spacing)) / 2
    start_y = Inches(4.6)
    
    for i, (title, desc) in enumerate(concepts):
        row = i // 4
        col = i % 4
        x = start_x + col * (card_width + spacing)
        y = start_y + row * (card_height + spacing)
        
        add_rounded_rect(slide, x, y, card_width, card_height)
        
        add_text_box(slide, title,
                     x + Inches(0.05), y + Inches(0.02),
                     card_width - Inches(0.1), Inches(0.3),
                     font_size=11, color=Colors.ACCENT, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(slide, desc,
                     x + Inches(0.05), y + Inches(0.3),
                     card_width - Inches(0.1), Inches(0.3),
                     font_size=9, color=Colors.MUTED, align=PP_ALIGN.CENTER)
    
    # ==================== SLIDE 11: THANK YOU ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    
    add_text_box(slide, "Thank You",
                 Inches(0.8), Inches(1.5), Inches(11), Inches(1.5),
                 font_size=60, color=Colors.ACCENT, bold=True, align=PP_ALIGN.CENTER)
    
    add_text_box(slide, "Questions & Discussion",
                 Inches(0.8), Inches(3.2), Inches(11), Inches(0.8),
                 font_size=28, color=Colors.SUBTEXT, align=PP_ALIGN.CENTER)
    
    add_text_box(slide, "Lagos State University — Department of Computer Science — CSC 226 — Group 6",
                 Inches(0.8), Inches(4.8), Inches(11), Inches(0.4),
                 font_size=14, color=Colors.MUTED, align=PP_ALIGN.CENTER)
    
    add_text_box(slide, "July 2026",
                 Inches(0.8), Inches(5.3), Inches(11), Inches(0.4),
                 font_size=14, color=Colors.MUTED, align=PP_ALIGN.CENTER)
    
    # Build command
    code14 = 'g++ -std=c++17 -I include main.cpp src/*.cpp -lws2_32  &&  ./certificate_generator'
    
    code_bg14 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.0), Inches(5.8),
        Inches(9.3), Inches(0.7)
    )
    code_bg14.fill.solid()
    code_bg14.fill.fore_color.rgb = Colors.CARD_BG
    code_bg14.line.color.rgb = Colors.BORDER
    
    add_text_box(slide, code14,
                 Inches(2.2), Inches(5.9),
                 Inches(8.9), Inches(0.5),
                 font_size=12, color=Colors.TEXT, font_name='Courier New', align=PP_ALIGN.CENTER)
    
    return prs

def main():
    """Main function to generate the presentation"""
    print("Generating LASU Certificate Generator Presentation...")
    
    try:
        prs = create_presentation()
        output_file = "LASU_Certificate_Generator.pptx"
        prs.save(output_file)
        print(f"✅ Presentation saved successfully as: {output_file}")
        print(f"📊 Total slides: {len(prs.slides)}")
        print("\n📝 Instructions:")
        print("1. Open the presentation in PowerPoint")
        print("2. On Slide 9, right-click the image placeholder and select 'Change Picture'")
        print("3. Select your 'preview.png' screenshot file")
        print("4. Adjust the image size if needed")
        print("\n🎨 The presentation uses the same color scheme as your HTML:")
        print(f"   - Background: #1a1817")
        print(f"   - Accent: #d97757")
        print(f"   - Text: #faf9f5")
        print(f"   - Cards: #2a2522 with #3a3532 borders")
        
    except Exception as e:
        print(f"❌ Error generating presentation: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    main()